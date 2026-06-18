import os
import re
import json
import uuid
import hashlib
import requests
import pandas as pd
from datetime import datetime
from functools import wraps
from flask import Flask, render_template, request, jsonify, send_file, session, Response, stream_with_context, redirect, url_for
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook, Workbook
from openpyxl.styles import Font
import xlrd
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from fpdf import FPDF

app = Flask(__name__)
app.secret_key = os.urandom(24)
UPLOAD = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')
USERS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'users.json')
SETTINGS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'settings.json')
ASSISTANTS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'assistants.json')
PERSONALIZATION_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'personalization.json')
app.config['UPLOAD_FOLDER'] = UPLOAD
app.config['OUTPUT_FOLDER'] = OUTPUT
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

ALLOWED = {'pdf', 'docx', 'xlsx', 'xls', 'pptx', 'txt', 'csv'}
OLLAMA_URL = os.environ.get('OLLAMA_URL', 'http://localhost:11434')
OLLAMA_MODEL = os.environ.get('OLLAMA_MODEL', 'qwen2.5:0.5b')
os.makedirs(UPLOAD, exist_ok=True)
os.makedirs(OUTPUT, exist_ok=True)

# ========== USER MANAGEMENT ==========
def load_users():
    if os.path.exists(USERS_FILE):
        with open(USERS_FILE, 'r') as f:
            return json.load(f)
    return {}

def save_users(users):
    with open(USERS_FILE, 'w') as f:
        json.dump(users, f, indent=2)

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if 'user' not in session:
            if request.is_json:
                return jsonify({"error": "Not authenticated"}), 401
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        data = request.get_json() if request.is_json else request.form
        username = data.get('username', '')
        password = data.get('password', '')
        users = load_users()
        if username in users and users[username] == hash_password(password):
            session['user'] = username
            if request.is_json:
                return jsonify({"success": True, "username": username})
            return redirect(url_for('index'))
        if request.is_json:
            return jsonify({"error": "Invalid credentials"}), 401
        return render_template('login.html', error='Invalid username or password')
    return render_template('login.html')

@app.route('/register', methods=['POST'])
def register():
    data = request.get_json() if request.is_json else request.form
    username = data.get('username', '')
    password = data.get('password', '')
    if not username or not password:
        if request.is_json:
            return jsonify({"error": "Username and password required"}), 400
        return render_template('login.html', error='Username and password required', show_register=True)
    users = load_users()
    if username in users:
        if request.is_json:
            return jsonify({"error": "Username already exists"}), 400
        return render_template('login.html', error='Username already exists', show_register=True)
    users[username] = hash_password(password)
    save_users(users)
    session['user'] = username
    if request.is_json:
        return jsonify({"success": True, "username": username})
    return redirect(url_for('index'))

@app.route('/logout')
def logout():
    session.pop('user', None)
    return redirect(url_for('login'))

# ========== SETTINGS MANAGEMENT ==========
def load_settings():
    if os.path.exists(SETTINGS_FILE):
        with open(SETTINGS_FILE, 'r') as f:
            return json.load(f)
    return {"theme": "dark", "language": "en", "notifications": True, "auto_save": True}

def save_settings(settings):
    with open(SETTINGS_FILE, 'w') as f:
        json.dump(settings, f, indent=2)

def load_assistants():
    if os.path.exists(ASSISTANTS_FILE):
        with open(ASSISTANTS_FILE, 'r') as f:
            return json.load(f)
    return []

def save_assistants(assistants):
    with open(ASSISTANTS_FILE, 'w') as f:
        json.dump(assistants, f, indent=2)

def load_personalization():
    if os.path.exists(PERSONALIZATION_FILE):
        with open(PERSONALIZATION_FILE, 'r') as f:
            return json.load(f)
    return {}

def save_personalization(data):
    with open(PERSONALIZATION_FILE, 'w') as f:
        json.dump(data, f, indent=2)

# ========== AI CORE ==========
def ollama_chat(messages, stream=False, model=None):
    try:
        resp = requests.post(f"{OLLAMA_URL}/api/chat",
            json={"model": model or OLLAMA_MODEL, "messages": messages, "stream": stream}, timeout=120)
        if stream:
            return resp
        data = resp.json()
        return data.get("message", {}).get("content", "")
    except Exception as e:
        return f"AI Error: {str(e)}"

def ollama_generate(prompt, stream=False, model=None):
    try:
        resp = requests.post(f"{OLLAMA_URL}/api/generate",
            json={"model": model or OLLAMA_MODEL, "prompt": prompt, "stream": stream}, timeout=120)
        if stream:
            return resp
        data = resp.json()
        return data.get("response", "")
    except Exception as e:
        return f"AI Error: {str(e)}"

def openai_chat(messages, api_key, model="gpt-3.5-turbo", stream=False):
    try:
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        data = {"model": model, "messages": messages, "stream": stream}
        resp = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=data, timeout=120)
        if stream:
            return resp
        result = resp.json()
        return result.get("choices", [{}])[0].get("message", {}).get("content", "")
    except Exception as e:
        return f"AI Error: {str(e)}"

def groq_chat(messages, api_key, model="llama3-8b-8192", stream=False):
    try:
        headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
        data = {"model": model, "messages": messages, "stream": stream}
        resp = requests.post("https://api.groq.com/openai/v1/chat/completions", headers=headers, json=data, timeout=120)
        if stream:
            return resp
        result = resp.json()
        return result.get("choices", [{}])[0].get("message", {}).get("content", "")
    except Exception as e:
        return f"AI Error: {str(e)}"

def gemini_chat(messages, api_key, model="gemini-1.5-flash", stream=False):
    try:
        contents = []
        for msg in messages:
            role = "user" if msg.get("role") == "user" else "model"
            if msg.get("role") == "system":
                contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
                contents.append({"role": "model", "parts": [{"text": "Understood."}]})
            else:
                contents.append({"role": role, "parts": [{"text": msg.get("content", "")}]})
        
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}"
        data = {"contents": contents}
        
        resp = requests.post(url, json=data, timeout=120)
        result = resp.json()
        
        if "candidates" in result and result["candidates"]:
            content = result["candidates"][0].get("content", {})
            parts = content.get("parts", [])
            if parts:
                return parts[0].get("text", "")
        
        error = result.get("error", {}).get("message", "Unknown error")
        return f"Gemini Error: {error}"
    except Exception as e:
        return f"Gemini Error: {str(e)}"

def gemini_stream(messages, api_key, model="gemini-1.5-flash"):
    try:
        contents = []
        for msg in messages:
            role = "user" if msg.get("role") == "user" else "model"
            if msg.get("role") == "system":
                contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
                contents.append({"role": "model", "parts": [{"text": "Understood."}]})
            else:
                contents.append({"role": role, "parts": [{"text": msg.get("content", "")}]})
        
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:streamGenerateContent?alt=sse&key={api_key}"
        data = {"contents": contents}
        
        resp = requests.post(url, json=data, timeout=120, stream=True)
        return resp
    except Exception as e:
        raise Exception(f"Gemini Error: {str(e)}")

def check_ollama():
    try:
        resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        models = resp.json().get("models", [])
        return {"status": "online", "models": [m["name"] for m in models], "current_model": OLLAMA_MODEL}
    except:
        return {"status": "offline", "models": [], "current_model": OLLAMA_MODEL}

# ========== FILE HELPERS ==========
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED

def get_file_type(filename):
    return filename.rsplit('.', 1)[1].lower() if '.' in filename else ''

# ========== EXTRACTION ==========
def extract_pdf(filepath):
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return {"text": text.strip(), "pages": len(reader.pages)}

def extract_docx(filepath):
    doc = Document(filepath)
    text = "\n".join([p.text for p in doc.paragraphs])
    tables = [[[cell.text for cell in row.cells] for row in t.rows] for t in doc.tables]
    return {"text": text.strip(), "tables": tables, "paragraphs": len(doc.paragraphs)}

def extract_xlsx(filepath):
    ext = get_file_type(filepath)
    if ext == 'xls':
        return extract_xls(filepath)
    wb = load_workbook(filepath, data_only=True)
    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        sheets[name] = [[str(c) if c is not None else "" for c in row] for row in ws.iter_rows(values_only=True)]
    return {"sheets": sheets, "sheet_names": wb.sheetnames}

def extract_xls(filepath):
    workbook = xlrd.open_workbook(filepath)
    sheets = {}
    sheet_names = []
    for sheet in workbook.sheets():
        sheet_names.append(sheet.name)
        rows = []
        for row_idx in range(sheet.nrows):
            row = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
            rows.append(row)
        sheets[sheet.name] = rows
    return {"sheets": sheets, "sheet_names": sheet_names}

def extract_pptx(filepath):
    prs = Presentation(filepath)
    slides = []
    for slide in prs.slides:
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text)
        slides.append(content)
    return {"slides": slides, "total_slides": len(prs.slides)}

def extract_csv(filepath):
    df = pd.read_csv(filepath)
    return {"data": df.to_dict(orient='records'), "columns": list(df.columns), "rows": len(df)}

def extract_text(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return {"text": f.read()}

def get_file_content(filepath):
    ext = get_file_type(filepath)
    if ext == 'pdf': return extract_pdf(filepath).get('text', '')
    elif ext == 'docx': return extract_docx(filepath).get('text', '')
    elif ext in ['xlsx', 'xls']:
        result = extract_xlsx(filepath)
        parts = [f"Sheet: {n}\n" + "\n".join(["\t".join(r) for r in rows[:100]]) for n, rows in result.get('sheets', {}).items()]
        return "\n\n".join(parts)
    elif ext == 'pptx':
        result = extract_pptx(filepath)
        return "\n\n".join([f"Slide {i+1}:\n" + "\n".join(s) for i, s in enumerate(result.get('slides', []))])
    elif ext == 'csv': return json.dumps(extract_csv(filepath).get('data', [])[:100], indent=2)
    elif ext == 'txt':
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f: return f.read()
    return ""

# ========== CONVERSION ==========
def smart_convert_pdf(structured, output_path):
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('DocTitle', parent=styles['Title'], fontSize=20, spaceAfter=12)
    heading_style = ParagraphStyle('DocHeading', parent=styles['Heading1'], fontSize=14, spaceAfter=8, textColor=colors.HexColor('#1e40af'))
    normal_style = ParagraphStyle('DocNormal', parent=styles['Normal'], fontSize=10, leading=14, spaceAfter=4)
    story = []
    if structured.get('title'):
        story.append(Paragraph(structured['title'], title_style))
        story.append(Spacer(1, 12))
    for section in structured.get('sections', []):
        if section.get('heading'):
            story.append(Paragraph(section['heading'], heading_style))
        for item in section.get('content', []):
            if item['type'] == 'text':
                story.append(Paragraph(item['value'], normal_style))
            elif item['type'] == 'table':
                tdata = item['value']
                if tdata and len(tdata) > 0:
                    t = Table(tdata)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2563eb')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                        ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db')),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8fafc')]),
                        ('TOPPADDING', (0, 0), (-1, -1), 4),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 10))
            elif item['type'] == 'list':
                for bullet in item['value']:
                    story.append(Paragraph(f"• {bullet}", normal_style))
    if not story:
        story.append(Paragraph(structured.get('text', 'No content'), normal_style))
    doc.build(story)

def smart_convert_docx(structured, output_path):
    doc = Document()
    if structured.get('title'):
        doc.add_heading(structured['title'], 0)
    for section in structured.get('sections', []):
        if section.get('heading'):
            doc.add_heading(section['heading'], level=1)
        for item in section.get('content', []):
            if item['type'] == 'text':
                doc.add_paragraph(item['value'])
            elif item['type'] == 'table':
                tdata = item['value']
                if tdata and len(tdata) > 0:
                    table = doc.add_table(rows=len(tdata), cols=len(tdata[0]) if tdata else 0)
                    table.style = 'Light Grid Accent 1'
                    for i, row in enumerate(tdata):
                        for j, cell in enumerate(row):
                            table.rows[i].cells[j].text = str(cell)
                    doc.add_paragraph()
            elif item['type'] == 'list':
                for bullet in item['value']:
                    doc.add_paragraph(bullet, style='List Bullet')
    if not any(s.get('content') for s in structured.get('sections', [])):
        doc.add_paragraph(structured.get('text', 'No content'))
    doc.save(output_path)

def smart_convert_xlsx(structured, output_path):
    wb = Workbook()
    ws = wb.active
    ws.title = structured.get('title', 'Data')[:31]
    all_tables = []
    for section in structured.get('sections', []):
        for item in section.get('content', []):
            if item['type'] == 'table':
                all_tables.append(item['value'])
    if all_tables:
        for ti, tdata in enumerate(all_tables):
            if ti > 0:
                ws = wb.create_sheet(title=f"Sheet{ti+1}")
            if tdata and len(tdata) > 0:
                for ri, row in enumerate(tdata):
                    for ci, cell in enumerate(row):
                        ws.cell(row=ri+1, column=ci+1, value=str(cell))
                for ci in range(len(tdata[0])):
                    ws.cell(row=1, column=ci+1).font = Font(bold=True)
    else:
        row = 1
        for section in structured.get('sections', []):
            if section.get('heading'):
                ws.cell(row=row, column=1, value=section['heading']).font = Font(bold=True, size=12)
                row += 1
            for item in section.get('content', []):
                if item['type'] == 'text':
                    ws.cell(row=row, column=1, value=item['value'])
                    row += 1
                elif item['type'] == 'list':
                    for bullet in item['value']:
                        ws.cell(row=row, column=1, value=bullet)
                        row += 1
        if row == 1:
            ws.cell(row=1, column=1, value=structured.get('text', 'No content'))
    for ws_item in wb.worksheets:
        for col in ws_item.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            ws_item.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)
    wb.save(output_path)

def smart_convert_pptx(structured, output_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = structured.get('title', 'Presentation')
    slide.placeholders[1].text = "Generated by FileArchitect"
    for section in structured.get('sections', []):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = section.get('heading', 'Section')
        texts = []
        tables = []
        for item in section.get('content', []):
            if item['type'] == 'text':
                texts.append(item['value'])
            elif item['type'] == 'list':
                texts.extend([f"• {b}" for b in item['value']])
            elif item['type'] == 'table':
                tables.append(item['value'])
        if texts:
            s.placeholders[1].text = "\n".join(texts[:8])
        elif tables and tables[0]:
            tdata = tables[0]
            rows, cols = len(tdata), len(tdata[0]) if tdata else 0
            tbl = s.shapes.add_table(rows, cols, 0.5, 2.0, 9.0, min(rows * 0.4, 4.0)).table
            for ri, row in enumerate(tdata):
                for ci, cell in enumerate(row):
                    tbl.cell(ri, ci).text = str(cell)
    if len(prs.slides) == 1:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Content"
        s.placeholders[1].text = structured.get('text', 'No content')[:500]
    prs.save(output_path)

def structure_from_file(filepath):
    ext = get_file_type(filepath)
    structured = {'title': os.path.splitext(os.path.basename(filepath))[0], 'sections': [], 'text': ''}
    if ext == 'pdf':
        result = extract_pdf(filepath)
        structured['text'] = result.get('text', '')
        structured['sections'] = parse_text_to_sections(result.get('text', ''))
    elif ext == 'docx':
        result = extract_docx(filepath)
        structured['text'] = result.get('text', '')
        if result.get('tables'):
            structured['sections'] = [{'heading': '', 'content': [{'type': 'table', 'value': t} for t in result['tables']]}]
        else:
            structured['sections'] = parse_text_to_sections(result.get('text', ''))
    elif ext in ['xlsx', 'xls']:
        result = extract_xlsx(filepath)
        for name, rows in result.get('sheets', {}).items():
            structured['sections'].append({'heading': name, 'content': [{'type': 'table', 'value': rows}]})
    elif ext == 'pptx':
        result = extract_pptx(filepath)
        for i, slide_content in enumerate(result.get('slides', [])):
            structured['sections'].append({'heading': f'Slide {i+1}', 'content': [{'type': 'text', 'value': '\n'.join(slide_content)}]})
    elif ext == 'csv':
        result = extract_csv(filepath)
        if result.get('data'):
            headers = result.get('columns', [])
            rows = [headers] + [[str(r.get(h, '')) for h in headers] for r in result['data']]
            structured['sections'] = [{'heading': '', 'content': [{'type': 'table', 'value': rows}]}]
    elif ext == 'txt':
        result = extract_text(filepath)
        structured['text'] = result.get('text', '')
        structured['sections'] = parse_text_to_sections(result.get('text', ''))
    return structured

def parse_text_to_sections(text):
    if not text:
        return []
    sections = []
    current = {'heading': '', 'content': []}
    lines = text.split('\n')
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.isupper() and len(stripped) > 3 and len(stripped) < 100:
            if current.get('content'):
                sections.append(current)
            current = {'heading': stripped, 'content': []}
        elif stripped.startswith('#'):
            if current.get('content'):
                sections.append(current)
            current = {'heading': stripped.lstrip('#').strip(), 'content': []}
        elif re.match(r'^\d+[\.\)]\s', stripped) or stripped.startswith('- ') or stripped.startswith('* '):
            if not current.get('content') or current['content'][-1].get('type') != 'list':
                current['content'].append({'type': 'list', 'value': []})
            current['content'][-1]['value'].append(stripped.lstrip('0123456789.)*- ').strip())
        elif '\t' in line or (stripped.count(',') >= 2 and len(stripped) < 500):
            cells = [c.strip() for c in re.split(r'[\t,]+', stripped)]
            if len(cells) >= 2:
                if not current.get('content') or current['content'][-1].get('type') != 'table':
                    current['content'].append({'type': 'table', 'value': []})
                current['content'][-1]['value'].append(cells)
        else:
            current['content'].append({'type': 'text', 'value': stripped})
    if current.get('content') or current.get('heading'):
        sections.append(current)
    return sections

def simple_text_to_pdf(text, output_path):
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    for line in text.split('\n'):
        if line.strip():
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 6))
    doc.build(story)

def simple_text_to_docx(text, output_path):
    doc = Document()
    for line in text.split('\n'):
        doc.add_paragraph(line)
    doc.save(output_path)

def simple_text_to_xlsx(text, output_path):
    wb = Workbook()
    ws = wb.active
    for i, line in enumerate(text.split('\n'), 1):
        ws.cell(row=i, column=1, value=line)
    wb.save(output_path)

def simple_text_to_pptx(text, output_path):
    prs = Presentation()
    lines = text.split('\n')
    for i in range(0, len(lines), 5):
        chunk = lines[i:i+5]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i // 5 + 1}"
        slide.placeholders[1].text = "\n".join(chunk)
    prs.save(output_path)

def data_to_pdf(data, output_path, title="Report"):
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles['Title']), Spacer(1, 20)]
    if data and isinstance(data[0], dict):
        headers = list(data[0].keys())
        table_data = [headers] + [[str(r.get(h, "")) for h in headers] for r in data]
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2563eb')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#d1d5db')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f9fafb')]),
        ]))
        story.append(table)
    doc.build(story)

def data_to_xlsx(data, output_path, sheet_name="Sheet1"):
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    if data and isinstance(data[0], dict):
        headers = list(data[0].keys())
        ws.append(headers)
        for row in data:
            ws.append([row.get(h, "") for h in headers])
    wb.save(output_path)

def data_to_docx(data, output_path, title="Report"):
    doc = Document()
    doc.add_heading(title, 0)
    for item in data:
        if isinstance(item, dict):
            for k, v in item.items():
                doc.add_paragraph(f"{k}: {v}", style='List Bullet')
        else:
            doc.add_paragraph(str(item))
    doc.save(output_path)

def data_to_pptx(data, output_path, title="Presentation"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    for i in range(0, len(data), 6):
        chunk = data[i:i+6]
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"Section {i // 6 + 1}"
        s.placeholders[1].text = "\n".join([str(x) for x in chunk])
    prs.save(output_path)

# ========== REPORT GENERATION ==========
def generate_student_report(data, output_path):
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=24, textColor=colors.HexColor('#1e40af'))
    story = [Paragraph("Student Academic Report", title_style), Spacer(1, 10)]
    for student in data.get('students', []):
        story.append(Paragraph(f"Student: {student.get('name', 'N/A')}", styles['Heading2']))
        story.append(Paragraph(f"Grade: {student.get('grade', 'N/A')}", styles['Normal']))
        marks = student.get('marks', {})
        if marks:
            table_data = [["Subject", "Marks", "Grade"]]
            for subj, mark in marks.items():
                if isinstance(mark, (int, float)):
                    g = "A" if mark >= 81 else "B" if mark >= 61 else "C" if mark >= 41 else "D" if mark >= 21 else "E"
                    table_data.append([subj, str(mark), g])
            table = Table(table_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e40af')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(table)
            story.append(Spacer(1, 15))
    doc.build(story)

def generate_certificate(data, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 28)
    pdf.cell(0, 40, "", ln=True)
    pdf.cell(0, 20, "Certificate of Achievement", ln=True, align="C")
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 20, "", ln=True)
    pdf.cell(0, 10, "This is to certify that", ln=True, align="C")
    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(0, 15, data.get('name', 'Student Name'), ln=True, align="C")
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 15, f"has successfully achieved {data.get('achievement', 'Academic Excellence')}", ln=True, align="C")
    pdf.cell(0, 10, f"Date: {data.get('date', datetime.now().strftime('%B %d, %Y'))}", ln=True, align="C")
    pdf.ln(30)
    pdf.set_font("Helvetica", "I", 12)
    pdf.cell(0, 10, "Authorized Signature", ln=True, align="C")
    pdf.output(output_path)

# ========== ANALYTICS ==========
def calculate_analytics(data):
    if not isinstance(data, list) or not data:
        return {}
    analytics = {"total_records": len(data)}
    num, cat = {}, {}
    for row in data:
        for k, v in row.items():
            if isinstance(v, (int, float)):
                num.setdefault(k, []).append(v)
            elif isinstance(v, str) and v:
                cat.setdefault(k, {})[v] = cat.get(k, {}).get(v, 0) + 1
    analytics["numeric_stats"] = {f: {"mean": round(sum(v)/len(v), 2), "min": min(v), "max": max(v), "sum": sum(v), "count": len(v)} for f, v in num.items()}
    analytics["categorical_stats"] = cat
    return analytics

# ========== ROUTES ==========
@app.route('/')
@login_required
def index():
    return render_template('index.html', user=session.get('user'))

@app.route('/api/ollama/status')
def ollama_status():
    return jsonify(check_ollama())

@app.route('/api/ollama/models')
def ollama_models():
    status = check_ollama()
    return jsonify({"models": status.get("models", [])})

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return jsonify({"error": "Invalid file"}), 400
    filename = secure_filename(file.filename)
    uid = str(uuid.uuid4())[:8]
    saved = f"{uid}_{filename}"
    filepath = os.path.join(UPLOAD, saved)
    file.save(filepath)
    session['last_upload'] = {'path': filepath, 'name': saved, 'original_name': filename, 'type': get_file_type(saved)}
    return jsonify({"success": True, "filename": saved, "original_name": filename, "file_type": get_file_type(saved)})

@app.route('/api/extract', methods=['POST'])
def extract_data():
    data = request.get_json()
    filename = data.get('filename') or session.get('last_upload', {}).get('name')
    if not filename:
        return jsonify({"error": "No file specified"}), 400
    filepath = os.path.join(UPLOAD, filename)
    if not os.path.exists(filepath):
        return jsonify({"error": "File not found"}), 404
    ext = get_file_type(filepath)
    try:
        if ext == 'pdf': result = extract_pdf(filepath)
        elif ext == 'docx': result = extract_docx(filepath)
        elif ext in ['xlsx', 'xls']: result = extract_xlsx(filepath)
        elif ext == 'pptx': result = extract_pptx(filepath)
        elif ext == 'csv': result = extract_csv(filepath)
        elif ext == 'txt': result = extract_text(filepath)
        else: return jsonify({"error": f"Unsupported: {ext}"}), 400
        result['file_type'] = ext
        result['filename'] = os.path.basename(filepath)
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/convert', methods=['POST'])
def convert_file():
    data = request.get_json()
    filename = data.get('filename')
    target = data.get('target_format')
    if not filename or not target:
        return jsonify({"error": "Missing parameters"}), 400
    filepath = os.path.join(UPLOAD, filename)
    if not os.path.exists(filepath):
        return jsonify({"error": "File not found"}), 404
    src = get_file_type(filepath)
    uid = str(uuid.uuid4())[:8]
    try:
        out_name = f"{uid}_converted.{target}"
        out_path = os.path.join(OUTPUT, out_name)
        if target == 'csv' and src in ['xlsx', 'xls']:
            if src == 'xls':
                wb_xls = xlrd.open_workbook(filepath)
                ws_xls = wb_xls.sheet_by_index(0)
                rows = []
                for row_idx in range(ws_xls.nrows):
                    rows.append([ws_xls.cell_value(row_idx, col_idx) for col_idx in range(ws_xls.ncols)])
                pd.DataFrame(rows[1:], columns=rows[0] if rows else []).to_csv(out_path, index=False)
            else:
                wb = load_workbook(filepath, data_only=True)
                ws = wb.active
                rows = list(ws.iter_rows(values_only=True))
                pd.DataFrame(rows[1:], columns=rows[0] if rows else []).to_csv(out_path, index=False)
        elif target == 'txt':
            content = get_file_content(filepath)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(content)
        else:
            structured = structure_from_file(filepath)
            if target == 'pdf': smart_convert_pdf(structured, out_path)
            elif target == 'docx': smart_convert_docx(structured, out_path)
            elif target == 'xlsx': smart_convert_xlsx(structured, out_path)
            elif target == 'pptx': smart_convert_pptx(structured, out_path)
            else: return jsonify({"error": f"Cannot convert {src} to {target}"}), 400
        return jsonify({"success": True, "output_filename": out_name, "message": f"Converted {src} to {target}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate', methods=['POST'])
def generate_report():
    data = request.get_json()
    report_type = data.get('report_type', 'pdf')
    template = data.get('template_type', 'general')
    content = data.get('content', {})
    uid = str(uuid.uuid4())[:8]
    try:
        if template == 'student_report':
            out_name = f"{uid}_student_report.pdf"
            generate_student_report(content if isinstance(content, dict) else {'text': str(content)}, os.path.join(OUTPUT, out_name))
        elif template == 'certificate':
            out_name = f"{uid}_certificate.pdf"
            generate_certificate(content if isinstance(content, dict) else {'name': str(content)}, os.path.join(OUTPUT, out_name))
        else:
            out_name = f"{uid}_report.{report_type}"
            out_path = os.path.join(OUTPUT, out_name)
            if isinstance(content, dict):
                text = content.get('text', content.get('data', ''))
            else:
                text = str(content)
            if isinstance(text, list): text = "\n".join([str(x) for x in text])
            if report_type == 'pdf': simple_text_to_pdf(str(text), out_path)
            elif report_type == 'docx': simple_text_to_docx(str(text), out_path)
            elif report_type == 'xlsx': simple_text_to_xlsx(str(text), out_path)
            elif report_type == 'pptx': simple_text_to_pptx(str(text), out_path)
        return jsonify({"success": True, "output_filename": out_name, "message": f"Generated {template}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/analytics', methods=['POST'])
def get_analytics():
    data = request.get_json()
    content = data.get('content', [])
    if not content: return jsonify({"error": "No data"}), 400
    try: return jsonify(calculate_analytics(content))
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/download/<filename>')
def download_file(filename):
    filepath = os.path.join(OUTPUT, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

# ========== SETTINGS ROUTES ==========
@app.route('/api/settings', methods=['GET'])
@login_required
def get_settings():
    return jsonify(load_settings())

@app.route('/api/settings', methods=['POST'])
@login_required
def update_settings():
    data = request.get_json()
    settings = load_settings()
    settings.update(data)
    save_settings(settings)
    return jsonify({"success": True, "settings": settings})

@app.route('/api/assistants', methods=['GET'])
@login_required
def get_assistants():
    return jsonify(load_assistants())

@app.route('/api/assistants', methods=['POST'])
@login_required
def add_assistant():
    data = request.get_json()
    assistants = load_assistants()
    assistant = {
        "id": str(uuid.uuid4())[:8],
        "name": data.get("name", "Assistant"),
        "provider": data.get("provider", "ollama"),
        "model": data.get("model", "qwen2.5:0.5b"),
        "api_key": data.get("api_key", ""),
        "api_url": data.get("api_url", ""),
        "system_prompt": data.get("system_prompt", "You are a helpful assistant."),
        "created_at": datetime.now().isoformat()
    }
    assistants.append(assistant)
    save_assistants(assistants)
    return jsonify({"success": True, "assistant": assistant})

@app.route('/api/assistants/<assistant_id>', methods=['DELETE'])
@login_required
def delete_assistant(assistant_id):
    assistants = load_assistants()
    assistants = [a for a in assistants if a.get("id") != assistant_id]
    save_assistants(assistants)
    return jsonify({"success": True})

@app.route('/api/assistants/<assistant_id>', methods=['PUT'])
@login_required
def update_assistant(assistant_id):
    data = request.get_json()
    assistants = load_assistants()
    for a in assistants:
        if a.get("id") == assistant_id:
            a.update({k: v for k, v in data.items() if k != "id"})
            save_assistants(assistants)
            return jsonify({"success": True, "assistant": a})
    return jsonify({"error": "Not found"}), 404

# ========== PERSONALIZATION ROUTES ==========
@app.route('/api/personalization', methods=['GET'])
@login_required
def get_personalization():
    return jsonify(load_personalization())

@app.route('/api/personalization', methods=['POST'])
@login_required
def update_personalization():
    data = request.get_json()
    save_personalization(data)
    return jsonify({"success": True})

# ========== SYSTEM DASHBOARD ==========
@app.route('/api/system', methods=['GET'])
@login_required
def get_system_info():
    import psutil
    cpu_percent = psutil.cpu_percent(interval=0.1)
    memory = psutil.virtual_memory()
    disk = psutil.disk_usage('/')
    
    uploads_size = 0
    if os.path.exists(UPLOAD):
        for f in os.listdir(UPLOAD):
            fp = os.path.join(UPLOAD, f)
            if os.path.isfile(fp):
                uploads_size += os.path.getsize(fp)
    
    models = []
    try:
        resp = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        for m in resp.json().get("models", []):
            models.append({"name": m["name"], "size": m.get("size", 0)})
    except:
        pass
    
    return jsonify({
        "cpu_percent": cpu_percent,
        "memory_total": memory.total,
        "memory_used": memory.used,
        "memory_percent": memory.percent,
        "disk_total": disk.total,
        "disk_used": disk.used,
        "disk_percent": disk.percent,
        "uploads_size": uploads_size,
        "models": models,
        "ollama_url": OLLAMA_URL,
        "current_model": OLLAMA_MODEL
    })

# ========== AI ROUTES ==========
@app.route('/api/ai/chat', methods=['POST'])
def ai_chat():
    data = request.get_json()
    messages = data.get('messages', [])
    filename = data.get('filename')
    assistant_id = data.get('assistant_id')
    request_model = data.get('model')
    
    file_content = ""
    if filename:
        fp = os.path.join(UPLOAD, filename)
        if os.path.exists(fp):
            file_content = get_file_content(fp)[:8000]
    
    system_msg = "You are FileArchitect AI. Help users analyze files, generate documents, and answer questions. Be helpful and concise."
    
    if file_content:
        system_msg += f"\n\nUploaded file content:\n---\n{file_content}\n---"
    
    full_messages = [{"role": "system", "content": system_msg}] + messages[-20:]
    
    provider = "ollama"
    model = request_model or OLLAMA_MODEL
    api_key = ""
    api_url = ""
    custom_prompt = ""
    
    if assistant_id:
        assistants = load_assistants()
        for a in assistants:
            if a.get("id") == assistant_id:
                provider = a.get("provider", "ollama")
                model = a.get("model", OLLAMA_MODEL)
                api_key = a.get("api_key", "")
                api_url = a.get("api_url", "")
                custom_prompt = a.get("system_prompt", "")
                if custom_prompt:
                    full_messages[0]["content"] = custom_prompt + "\n\n" + (file_content if file_content else "")
                break
    
    def generate():
        try:
            if provider == "openai":
                resp = openai_chat(full_messages, api_key, model, stream=True)
                for line in resp.iter_lines():
                    if line:
                        try:
                            chunk = json.loads(line)
                            token = chunk.get("choices", [{}])[0].get("delta", {}).get("content", "")
                            if token:
                                yield f"data: {json.dumps({'token': token})}\n\n"
                        except json.JSONDecodeError:
                            continue
            elif provider == "groq":
                resp = groq_chat(full_messages, api_key, model, stream=True)
                for line in resp.iter_lines():
                    if line:
                        try:
                            chunk = json.loads(line)
                            token = chunk.get("choices", [{}])[0].get("delta", {}).get("content", "")
                            if token:
                                yield f"data: {json.dumps({'token': token})}\n\n"
                        except json.JSONDecodeError:
                            continue
            elif provider == "gemini":
                resp = gemini_stream(full_messages, api_key, model)
                for line in resp.iter_lines():
                    if line:
                        try:
                            line_str = line.decode('utf-8') if isinstance(line, bytes) else line
                            if line_str.startswith('data: '):
                                chunk = json.loads(line_str[6:])
                                candidates = chunk.get("candidates", [])
                                if candidates:
                                    parts = candidates[0].get("content", {}).get("parts", [])
                                    if parts:
                                        token = parts[0].get("text", "")
                                        if token:
                                            yield f"data: {json.dumps({'token': token})}\n\n"
                        except (json.JSONDecodeError, UnicodeDecodeError):
                            continue
            else:
                resp = ollama_chat(full_messages, stream=True, model=model)
                for line in resp.iter_lines():
                    if line:
                        try:
                            chunk = json.loads(line)
                            token = chunk.get("message", {}).get("content", "")
                            if token:
                                yield f"data: {json.dumps({'token': token})}\n\n"
                        except json.JSONDecodeError:
                            continue
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"
        yield "data: [DONE]\n\n"
    
    return Response(stream_with_context(generate()), content_type='text/event-stream')

@app.route('/api/ai/summarize', methods=['POST'])
def ai_summarize():
    data = request.get_json()
    filename = data.get('filename')
    if not filename: return jsonify({"error": "No file"}), 400
    fp = os.path.join(UPLOAD, filename)
    if not os.path.exists(fp): return jsonify({"error": "Not found"}), 404
    content = get_file_content(fp)[:6000]
    result = ollama_generate(f"Summarize this document concisely with key points:\n\n{content}")
    return jsonify({"summary": result})

@app.route('/api/ai/analyze', methods=['POST'])
def ai_analyze():
    data = request.get_json()
    filename = data.get('filename')
    question = data.get('question', 'Analyze this data and provide insights')
    if not filename: return jsonify({"error": "No file"}), 400
    fp = os.path.join(UPLOAD, filename)
    if not os.path.exists(fp): return jsonify({"error": "Not found"}), 404
    content = get_file_content(fp)[:6000]
    result = ollama_generate(f"Analyze this data and answer the question.\n\nData:\n{content}\n\nQuestion: {question}")
    return jsonify({"analysis": result})

@app.route('/api/ai/generate', methods=['POST'])
def ai_generate():
    data = request.get_json()
    prompt = data.get('prompt', '')
    fmt = data.get('output_format', 'txt')
    if not prompt: return jsonify({"error": "No prompt"}), 400
    result = ollama_generate(f"Generate content for this request. Return ONLY the content.\n\nRequest: {prompt}\n\nFormat for {fmt} file.")
    uid = str(uuid.uuid4())[:8]
    out_name = f"{uid}_ai_generated.{fmt}"
    out_path = os.path.join(OUTPUT, out_name)
    try:
        if fmt == 'pdf': simple_text_to_pdf(result, out_path)
        elif fmt == 'docx': simple_text_to_docx(result, out_path)
        elif fmt == 'pptx': simple_text_to_pptx(result, out_path)
        elif fmt == 'xlsx': simple_text_to_xlsx(result, out_path)
        else:
            with open(out_path, 'w', encoding='utf-8') as f: f.write(result)
        return jsonify({"success": True, "output_filename": out_name, "content": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ai/extract_smart', methods=['POST'])
def ai_extract_smart():
    data = request.get_json()
    filename = data.get('filename')
    extract_type = data.get('extract_type', 'entities')
    if not filename: return jsonify({"error": "No file"}), 400
    fp = os.path.join(UPLOAD, filename)
    if not os.path.exists(fp): return jsonify({"error": "Not found"}), 404
    content = get_file_content(fp)[:6000]
    prompts = {
        "entities": f"Extract all key entities (names, dates, locations, orgs, amounts) as a structured list:\n\n{content}",
        "table": f"Extract tabular data and return as structured table:\n\n{content}",
        "key_values": f"Extract key-value pairs as JSON:\n\n{content}",
        "contacts": f"Extract contact info (names, emails, phones, addresses):\n\n{content}",
        "financial": f"Extract financial data (amounts, dates, descriptions):\n\n{content}"
    }
    result = ollama_generate(prompts.get(extract_type, prompts["entities"]))
    return jsonify({"extracted": result, "type": extract_type})

# ========== EXTENSIONS ==========
from extensions import registry

@app.route('/api/extensions', methods=['GET'])
@login_required
def get_extensions():
    return jsonify(registry.get_all())

@app.route('/api/extensions/<name>/enable', methods=['POST'])
@login_required
def enable_extension(name):
    if registry.enable(name):
        return jsonify({"success": True})
    return jsonify({"error": "Extension not found"}), 404

@app.route('/api/extensions/<name>/disable', methods=['POST'])
@login_required
def disable_extension(name):
    if registry.disable(name):
        return jsonify({"success": True})
    return jsonify({"error": "Extension not found"}), 404

@app.route('/api/extensions/<name>/config', methods=['POST'])
@login_required
def update_extension_config(name):
    data = request.get_json()
    if registry.update_config(name, data):
        return jsonify({"success": True})
    return jsonify({"error": "Extension not found"}), 404

@app.route('/api/extensions/reload', methods=['POST'])
@login_required
def reload_extensions():
    registry.reload()
    return jsonify({"success": True})

@app.route('/api/extensions/tools', methods=['GET'])
@login_required
def get_extension_tools():
    return jsonify(registry.get_chat_tools())

@app.route('/api/extensions/tools/<tool_name>', methods=['POST'])
@login_required
def execute_extension_tool(tool_name):
    data = request.get_json()
    result = registry.execute_tool(tool_name, data)
    return jsonify(result)

# Initialize extensions after all routes are defined
registry.init_app(app)

if __name__ == '__main__':
    import socket
    import subprocess
    try:
        result = subprocess.run(['powershell', '-Command', 'Get-NetIPAddress -AddressFamily IPv4 | Where-Object {$_.IPAddress -notlike "127.*" -and $_.IPAddress -notlike "169.254.*"} | Select-Object -First 1 -ExpandProperty IPAddress'], capture_output=True, text=True, timeout=5)
        local_ip = result.stdout.strip()
    except:
        local_ip = "localhost"
    if not local_ip:
        local_ip = "localhost"
    print(f"\n{'='*50}")
    print(f"  FILEARCHITECT // NEURAL INTERFACE")
    print(f"{'='*50}")
    print(f"  LOCAL:   http://localhost:5000")
    print(f"  NETWORK: http://{local_ip}:5000")
    print(f"{'='*50}")
    print(f"  Share NETWORK URL with others on same WiFi")
    print(f"  Login: admin / admin123")
    print(f"{'='*50}\n")
    app.run(debug=True, host='0.0.0.0', port=5000)
